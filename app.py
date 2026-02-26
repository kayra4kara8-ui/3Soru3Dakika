"""
POLCAST | Eczacı Elif Aracıoğlu | Video Stüdyo v16.0
──────────────────────────────────────────────────────────────
• Marka rengi: Eczacı yeşili #34A883, Cormorant Garamond + DM Sans
• Ses temizleme: highpass(80Hz) + afftdn(gürültü azaltma)
• v16: Profesyonel Medikal Jingle Motoru
  - Açılış: EKG bip zinciri + steril pad + diapason + kristal arpej
  - Kapanış: Yavaşlayan EKG + inen arpej + medikal onay çifti
• Görsel: Animasyonlu EKG çizgisi + nabız halkası + köşe dekor
• Slayt içeriği bantların DIŞINDA — başlıklar/içerik tam görünür
"""
import streamlit as st
import io, os, math, base64, tempfile, subprocess, shutil, json, time
import numpy as np

# ── ffmpeg ────────────────────────────────────────────────────────────────────
def _get_ffmpeg():
    try:
        import imageio_ffmpeg
        exe = imageio_ffmpeg.get_ffmpeg_exe()
        if exe and os.path.exists(exe):
            return exe
    except Exception:
        pass
    found = shutil.which("ffmpeg")
    if found:
        return found
    for p in ["/usr/bin/ffmpeg", "/usr/local/bin/ffmpeg", "/opt/homebrew/bin/ffmpeg"]:
        if os.path.exists(p):
            return p
    return None

FFMPEG = _get_ffmpeg()
if FFMPEG:
    os.environ["IMAGEIO_FFMPEG_EXE"] = FFMPEG

try:
    from PIL import Image, ImageDraw, ImageFont
    PIL_OK = True
except ImportError:
    PIL_OK = False
try:
    import imageio  # noqa
    IMAGEIO_OK = True
except ImportError:
    IMAGEIO_OK = False
try:
    from pptx import Presentation
    PPTX_OK = True
except ImportError:
    PPTX_OK = False

LO_BIN  = "/usr/bin/libreoffice" if os.path.exists("/usr/bin/libreoffice") else "/usr/bin/soffice"
LO_OK   = os.path.exists(LO_BIN)
PPM_OK  = os.path.exists("/usr/bin/pdftoppm")
FFMPEG_OK = FFMPEG is not None

VIDEO_W   = 1280
VIDEO_H   = 720
VIDEO_FPS = 24
TOP_BAR   = 58
BOT_BAR   = 48
SLIDE_AREA_H = VIDEO_H - TOP_BAR - BOT_BAR
SLIDE_AREA_W = VIDEO_W
BRAND_RGB = (52, 168, 131)
BRAND_HEX = "34A883"
JINGLE_DUR = 5.0

FONT_PATHS = [
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    "/usr/share/fonts/truetype/ubuntu/Ubuntu-B.ttf",
]

PALETTE = [
    {"hex": BRAND_HEX, "rgb": BRAND_RGB,     "emoji": "💊"},
    {"hex": "4C9FCA",  "rgb": (76,159,202),  "emoji": "👩‍💼"},
    {"hex": "C9A84C",  "rgb": (201,168,76),  "emoji": "🎤"},
    {"hex": "E07B7B",  "rgb": (195,90,90),   "emoji": "🎙️"},
    {"hex": "B57FCC",  "rgb": (155,105,195), "emoji": "💬"},
    {"hex": "7EC8C8",  "rgb": (80,178,178),  "emoji": "📢"},
    {"hex": "F0A060",  "rgb": (220,140,70),  "emoji": "🗣️"},
    {"hex": "88BBEE",  "rgb": (100,160,220), "emoji": "👤"},
]
DEFAULT_CHARACTERS = [{"name": "Elif Aracıoğlu", "role": "Eczacı", **PALETTE[0]}]

# ═════════════════════════════════════════════════════════════════════════════
# YARDIMCILAR
# ═════════════════════════════════════════════════════════════════════════════
def _font(size):
    if not PIL_OK: return None
    for p in FONT_PATHS:
        if os.path.exists(p):
            try: return ImageFont.truetype(p, size)
            except: pass
    return ImageFont.load_default()

def _run(cmd, timeout=900, step_name=""):
    try:
        r = subprocess.run(cmd, capture_output=True, text=True, timeout=timeout)
        if r.returncode != 0:
            raise RuntimeError(
                f"[{step_name}] Kod {r.returncode}\nCMD: {' '.join(str(c) for c in cmd)}\n"
                f"STDERR: {r.stderr[-800:] if r.stderr else '(yok)'}")
        return r
    except subprocess.TimeoutExpired:
        raise RuntimeError(f"[{step_name}] Zaman aşımı {timeout}s")

def _ffprobe_path():
    if FFMPEG:
        fp = FFMPEG.replace("ffmpeg", "ffprobe")
        if os.path.exists(fp): return fp
    found = shutil.which("ffprobe")
    if found: return found
    for p in ["/usr/bin/ffprobe", "/usr/local/bin/ffprobe"]:
        if os.path.exists(p): return p
    return "ffprobe"

def audio_duration_ffprobe(path: str) -> float:
    try:
        r = subprocess.run(
            [_ffprobe_path(), "-v", "error", "-show_entries", "format=duration",
             "-of", "json", path],
            capture_output=True, text=True, timeout=30)
        return float(json.loads(r.stdout)["format"]["duration"])
    except:
        size = os.path.getsize(path) if os.path.exists(path) else 0
        return max(1.0, size / 16000)

def audio_duration_sec_bytes(data: bytes) -> float:
    if not data: return 3.0
    tmp = tempfile.mktemp(suffix=".audio")
    try:
        with open(tmp, "wb") as f: f.write(data)
        return audio_duration_ffprobe(tmp)
    finally:
        try: os.unlink(tmp)
        except: pass

CLEAN_AF = (
    "highpass=f=80,"
    "afftdn=nf=-20,"
    "aresample=async=1:min_hard_comp=0.1:first_pts=0"
)

def clean_audio(inp: str, out: str, step: str = "Ses temizleme"):
    _run([FFMPEG, "-y", "-i", inp, "-af", CLEAN_AF,
          "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2", out],
         timeout=180, step_name=step)

# ═════════════════════════════════════════════════════════════════════════════
# PROFESYONEl MEDİKAL JİNGLE MOTORU v2.0
# ─────────────────────────────────────────────────────────────────────────────
# Açılış (open) — 5 saniye:
#   • EKG bip zinciri — 880 Hz klinik monitor sesi, P dalgası simülasyonu
#   • Steril sine pad — C2+G2+E3 major akoru, chorus ile yumuşatılmış
#   • Medikal whoosh — bandpass brown noise (200–2200 Hz sweep)
#   • Diapason — 440 Hz A referans tonu
#   • Kristal arpej — C5→E5→G5→B5→C6 pentatonik yükseliş + harmonik
#
# Kapanış (close) — 5 saniye:
#   • Yavaşlayan EKG — artan aralıklar (72→60 BPM hissi)
#   • Pad söner — aynı major akoru, yavaşça fadeout
#   • İnen kristal arpej — C6→B5→G5→E5→C5
#   • Medikal onay çifti — 1320+1760 Hz (tamamlandı sinyali)
#   • Son nefes — alçak brown noise, kapanış hissi
# ═════════════════════════════════════════════════════════════════════════════

def make_jingle(work_dir: str, kind: str = "open") -> str:
    """
    Profesyonel medikal jingle — tamamen ffmpeg lavfi sentezi, dış dosya yok.
    """
    out_path = os.path.join(work_dir, f"jingle_{kind}.aac")

    # Yardımcı: lavfi filter → AAC
    def _af(filt: str, dur: float, fname: str) -> str:
        p = os.path.join(work_dir, fname)
        subprocess.run(
            [FFMPEG, "-y", "-f", "lavfi", "-i", filt,
             "-t", str(dur), "-c:a", "aac", "-b:a", "128k",
             "-ar", "44100", "-ac", "2", p],
            capture_output=True, timeout=30
        )
        return p

    # EKG bip — 880 Hz, 60 ms klinik monitor sesi
    def _ekg(tag: str, vol: float = 0.55) -> str:
        return _af(
            f"sine=frequency=880:duration=0.06,"
            f"afade=t=in:st=0:d=0.004,"
            f"afade=t=out:st=0.025:d=0.035,"
            f"volume={vol}",
            0.06, f"ekg_{tag}.aac"
        )

    # P dalgası simülasyonu — 660 Hz, 30 ms
    def _pw(tag: str) -> str:
        return _af(
            "sine=frequency=660:duration=0.03,"
            "afade=t=in:st=0:d=0.003,"
            "afade=t=out:st=0.012:d=0.018,"
            "volume=0.18",
            0.03, f"pw_{tag}.aac"
        )

    # Ortak katmanlar
    whoosh = _af(
        "anoisesrc=color=brown:duration=5,"
        "highpass=f=200,"
        "lowpass=f=2200,"
        "afade=t=in:st=0:d=1.8,"
        "afade=t=out:st=3.2:d=1.8,"
        "volume=0.09",
        5.0, f"whoosh_{kind}.aac"
    )
    pad_c = _af(
        "sine=frequency=130.8:duration=5,"
        "chorus=0.4:0.8:45|55:0.25|0.25:0.18|0.22:2|1.5,"
        "volume=0.19",
        5.0, f"pad_c_{kind}.aac"
    )
    pad_g = _af(
        "sine=frequency=196.0:duration=5,"
        "chorus=0.4:0.8:50|60:0.25|0.25:0.18|0.22:2|1.5,"
        "volume=0.15",
        5.0, f"pad_g_{kind}.aac"
    )
    pad_e = _af(
        "sine=frequency=329.6:duration=5,"
        "chorus=0.4:0.8:40|50:0.25|0.25:0.18|0.22:2|1.5,"
        "volume=0.12",
        5.0, f"pad_e_{kind}.aac"
    )
    diapason = _af(
        "sine=frequency=440:duration=0.6,"
        "afade=t=in:st=0:d=0.04,"
        "afade=t=out:st=0.35:d=0.25,"
        "volume=0.28",
        0.6, f"diapason_{kind}.aac"
    )

    if kind == "open":
        # EKG ritmi: 72 BPM — P dalgası + QRS × 4
        e1, e2, e3, e4 = _ekg("o1", 0.42), _ekg("o2", 0.48), _ekg("o3", 0.54), _ekg("o4", 0.60)
        p1, p2, p3     = _pw("o1"), _pw("o2"), _pw("o3")

        # Kristal yükselen arpej: C5 → E5 → G5 → B5 → C6
        a1 = _af("sine=frequency=523.3:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.32", 0.35, "ao1.aac")
        a2 = _af("sine=frequency=659.3:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.32", 0.35, "ao2.aac")
        a3 = _af("sine=frequency=784.0:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.32", 0.35, "ao3.aac")
        a4 = _af("sine=frequency=987.8:duration=0.38,afade=t=in:st=0:d=0.015,afade=t=out:st=0.25:d=0.13,volume=0.34", 0.38, "ao4.aac")
        a5 = _af("sine=frequency=1046.5:duration=1.4,afade=t=in:st=0:d=0.05,afade=t=out:st=0.9:d=0.5,volume=0.36",   1.4,  "ao5.aac")
        harm = _af(
            "sine=frequency=2093:duration=0.8,"
            "afade=t=in:st=0:d=0.06,"
            "afade=t=out:st=0.5:d=0.3,"
            "volume=0.10",
            0.8, "ao_harm.aac"
        )

        parts = [
            (whoosh,   0),
            (pad_c,    200),   (pad_g,  200),  (pad_e, 280),
            (p1,        50),   (e1,      95),
            (p2,       928),   (e2,     973),
            (p3,      1761),   (e3,    1806),
                               (e4,    2639),
            (diapason, 2450),
            (a1,       3000),
            (a2,       3360),
            (a3,       3720),
            (a4,       4080),
            (a5,       4460),
            (harm,     4480),
        ]
        fade_out_at = 4.32

    else:
        # EKG yavaşlıyor: 833 → 950 → 1100 ms aralık
        e1, e2, e3 = _ekg("c1", 0.55), _ekg("c2", 0.50), _ekg("c3", 0.42)

        # Medikal çift onay bip
        conf1 = _af(
            "sine=frequency=1320:duration=0.12,"
            "afade=t=in:st=0:d=0.008,"
            "afade=t=out:st=0.07:d=0.05,"
            "volume=0.42",
            0.12, "conf1.aac"
        )
        conf2 = _af(
            "sine=frequency=1760:duration=0.16,"
            "afade=t=in:st=0:d=0.008,"
            "afade=t=out:st=0.09:d=0.07,"
            "volume=0.38",
            0.16, "conf2.aac"
        )

        # İnen kristal arpej: C6 → B5 → G5 → E5 → C5
        c1 = _af("sine=frequency=1046.5:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.34", 0.35, "ac1.aac")
        c2 = _af("sine=frequency=987.8:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.32",  0.35, "ac2.aac")
        c3 = _af("sine=frequency=784.0:duration=0.35,afade=t=in:st=0:d=0.015,afade=t=out:st=0.22:d=0.13,volume=0.30",  0.35, "ac3.aac")
        c4 = _af("sine=frequency=659.3:duration=0.38,afade=t=in:st=0:d=0.015,afade=t=out:st=0.25:d=0.13,volume=0.28",  0.38, "ac4.aac")
        c5 = _af("sine=frequency=523.3:duration=1.8,afade=t=in:st=0:d=0.07,afade=t=out:st=1.1:d=0.7,volume=0.32",      1.8,  "ac5.aac")
        breath = _af(
            "anoisesrc=color=brown:duration=1.2,"
            "highpass=f=80,"
            "lowpass=f=400,"
            "afade=t=in:st=0:d=0.4,"
            "afade=t=out:st=0.6:d=0.6,"
            "volume=0.06",
            1.2, "breath_c.aac"
        )

        parts = [
            (pad_c,    0),    (pad_g,    0),    (pad_e,   60),
            (whoosh,   0),
            (e1,       200),
            (e2,      1033),
            (e3,      1983),
            (diapason, 2200),
            (c1,       2600),
            (c2,       2960),
            (c3,       3320),
            (c4,       3680),
            (c5,       4040),
            (conf1,    4600),
            (conf2,    4740),
            (breath,   3800),
        ]
        fade_out_at = 4.50

    # Mix
    total_dur = 5.0
    fo_dur    = total_dur - fade_out_at
    in_args   = []
    for path, _ in parts:
        in_args += ["-i", path]

    n  = len(parts)
    fp = []
    for i, (_, delay_ms) in enumerate(parts):
        fp.append(f"[{i}]adelay={delay_ms}|{delay_ms}[d{i}]")
    fp.append(
        "".join(f"[d{i}]" for i in range(n)) +
        f"amix=inputs={n}:normalize=0[mix]"
    )
    fp.append(
        f"[mix]"
        f"afade=t=in:st=0:d=0.12,"
        f"afade=t=out:st={fade_out_at}:d={fo_dur},"
        f"volume=3.2"
        f"[out]"
    )

    subprocess.run(
        [FFMPEG, "-y", *in_args,
         "-filter_complex", ";".join(fp),
         "-map", "[out]",
         "-t", str(total_dur),
         "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2",
         out_path],
        capture_output=True, timeout=60
    )

    for path, _ in parts:
        try: os.unlink(path)
        except: pass

    return out_path


# ═════════════════════════════════════════════════════════════════════════════
# AÇILIŞ / KAPANIŞ KARESI RENDER — Medikal Temalı Animasyon v2.0
# ─────────────────────────────────────────────────────────────────────────────
# • Animasyonlu kayan EKG çizgisi (QRS kompleksi)
# • Aktif ışıma noktası
# • Çift nabız halkası
# • Pulse ile nefes alan medikal artı sembolü
# • Köşe dekorasyonları (steril çerçeve)
# • Süsleme nokta dizisi
# ═════════════════════════════════════════════════════════════════════════════

def render_intro_frame(t: float, kind: str, color=BRAND_RGB) -> np.ndarray:
    """
    t: 0.0–1.0 arası animasyon zamanı
    kind: 'open' veya 'close'
    """
    frame = Image.new("RGB", (VIDEO_W, VIDEO_H), (3, 6, 5))
    draw  = ImageDraw.Draw(frame, "RGBA")

    fn40 = _font(40)
    fn24 = _font(24)
    fn15 = _font(15)
    fn12 = _font(12)

    cx, cy = VIDEO_W // 2, VIDEO_H // 2

    # Arka plan radyal gradyan
    pulse_bg = 0.4 + 0.6 * abs(math.sin(t * math.pi))
    for r in range(380, 0, -35):
        alpha = int(18 * (1 - r / 380) * pulse_bg)
        draw.ellipse(
            [cx - r, cy - r, cx + r, cy + r],
            fill=(*color, max(0, min(255, alpha)))
        )

    # Animasyonlu kayan EKG çizgisi
    ekg_y  = cy + 128
    ekg_w  = 680
    ekg_x0 = cx - ekg_w // 2
    speed  = 1.6 if kind == "open" else -1.2
    offset = int(abs(t * ekg_w * speed)) % ekg_w

    def ekg_val(x_norm: float) -> float:
        x = (x_norm * 4.0) % 1.0
        if x < 0.20:
            return math.sin(x * math.pi / 0.20) * 7
        elif x < 0.30:
            return 0.0
        elif x < 0.34:
            return -30 * (x - 0.30) / 0.04
        elif x < 0.38:
            return -30 + 85 * (x - 0.34) / 0.04
        elif x < 0.43:
            return 55 - 75 * (x - 0.38) / 0.05
        elif x < 0.50:
            return -20 + 20 * (x - 0.43) / 0.07
        elif x < 0.68:
            return math.sin((x - 0.50) * math.pi / 0.18) * 13
        else:
            return 0.0

    prev_px = prev_py = None
    for px in range(ekg_x0, ekg_x0 + ekg_w, 2):
        local     = px - ekg_x0
        x_shifted = ((local + offset) % ekg_w) / ekg_w
        val = ekg_val(x_shifted)
        py  = int(ekg_y - val)
        if prev_px is not None:
            alpha = min(255, 120 + int(abs(val) * 2.5))
            draw.line([prev_px, prev_py, px, py],
                      fill=(*color, alpha), width=2)
        prev_px, prev_py = px, py

    # Aktif EKG ışıma noktası
    head_x_norm = ((offset + ekg_w // 3) % ekg_w) / ekg_w
    head_val    = ekg_val(head_x_norm)
    head_px     = ekg_x0 + ekg_w // 3
    head_py     = int(ekg_y - head_val)
    glow_r      = 5 + int(abs(head_val) * 0.08)
    draw.ellipse([head_px - glow_r, head_py - glow_r,
                  head_px + glow_r, head_py + glow_r],
                 fill=(*color, 240))
    draw.ellipse([head_px - 2, head_py - 2, head_px + 2, head_py + 2],
                 fill=(255, 255, 255, 200))

    # Üç animasyonlu halka
    r1 = int(145 + 18 * math.sin(t * math.pi * 2))
    r2 = int(168 + 12 * math.sin(t * math.pi * 2 + 1.1))
    r3 = int(188 + 8  * math.sin(t * math.pi * 1.5 + 0.5))
    draw.ellipse([cx - r3, cy - r3, cx + r3, cy + r3], outline=(*color, 18), width=1)
    draw.ellipse([cx - r2, cy - r2, cx + r2, cy + r2], outline=(*color, 38), width=1)
    draw.ellipse([cx - r1, cy - r1, cx + r1, cy + r1], outline=(*color, 85), width=2)

    # İç daire
    draw.ellipse([cx - 108, cy - 108, cx + 108, cy + 108],
                 fill=(*color, 22), outline=(*color, 115), width=3)

    # Medikal artı — pulse ile nefes alıyor
    pulse_scale = 0.86 + 0.14 * abs(math.sin(t * math.pi * 4))
    cs = int(32 * pulse_scale)
    cw = int(10 * pulse_scale)
    draw.rectangle([cx - cw//2 + 1, cy - cs + 1, cx + cw//2 + 1, cy + cs + 1],
                   fill=(*color, 40))
    draw.rectangle([cx - cs + 1, cy - cw//2 + 1, cx + cs + 1, cy + cw//2 + 1],
                   fill=(*color, 40))
    draw.rectangle([cx - cw//2, cy - cs, cx + cw//2, cy + cs],
                   fill=(*color, 215))
    draw.rectangle([cx - cs, cy - cw//2, cx + cs, cy + cw//2],
                   fill=(*color, 215))

    # Fade alpha
    if kind == "open":
        fade = min(1.0, t * 2.8)
    else:
        fade = max(0.0, 1.0 - t * 2.2)
    ta = int(255 * fade)

    # POLCAST başlığı
    title = "POLCAST"
    try:
        tw = draw.textlength(title, font=fn40)
    except:
        tw = len(title) * 22
    draw.text((cx - tw // 2, cy - 102), title, font=fn40, fill=(*color, ta))

    # Animasyonlu alt çizgi
    if kind == "open":
        lw = int(tw * min(1.0, t * 2.5))
    else:
        lw = int(tw * max(0.0, 1.0 - t * 2.0))
    draw.rectangle([cx - lw // 2, cy - 56, cx + lw // 2, cy - 53],
                   fill=(*color, ta))

    # EKG etiketi
    ekg_label = "♥  EKG" if kind == "open" else "♥  BYE"
    try:
        elw = draw.textlength(ekg_label, font=fn12)
    except:
        elw = len(ekg_label) * 7
    draw.text((cx - elw // 2, cy - 42), ekg_label,
              font=fn12, fill=(*color, int(ta * 0.50)))

    # İsim
    name = "Eczacı Elif Aracıoğlu"
    try:
        nw = draw.textlength(name, font=fn24)
    except:
        nw = len(name) * 13
    draw.text((cx - nw // 2, cy + 132), name,
              font=fn24, fill=(195, 235, 215, ta))

    # Açılış / kapanış etiketi
    label = "HOŞ GELDİNİZ" if kind == "open" else "GÖRÜŞMEK ÜZERE"
    try:
        llw = draw.textlength(label, font=fn15)
    except:
        llw = len(label) * 9
    draw.text((cx - llw // 2, cy + 170), label,
              font=fn15, fill=(*color, int(ta * 0.62)))

    # Süsleme nokta dizisi
    for i, bx in enumerate([cx - 220, cx - 110, cx, cx + 110, cx + 220]):
        dot_alpha = int(ta * (0.25 + 0.75 * abs(math.sin(t * math.pi * 3 + i * 0.9))))
        draw.ellipse([bx - 3, cy + 200, bx + 3, cy + 206],
                     fill=(*color, dot_alpha))

    # Üst / alt şerit
    draw.rectangle([0, 0, VIDEO_W, 4], fill=(*color, 200))
    draw.rectangle([0, VIDEO_H - 4, VIDEO_W, VIDEO_H], fill=(*color, 200))

    # Köşe dekorasyonları
    corner = 24
    ca     = int(ta * 0.45)
    draw.rectangle([0, 0, corner, 2],                       fill=(*color, ca))
    draw.rectangle([0, 0, 2, corner],                       fill=(*color, ca))
    draw.rectangle([VIDEO_W - corner, 0, VIDEO_W, 2],       fill=(*color, ca))
    draw.rectangle([VIDEO_W - 2, 0, VIDEO_W, corner],       fill=(*color, ca))
    draw.rectangle([0, VIDEO_H - 2, corner, VIDEO_H],       fill=(*color, ca))
    draw.rectangle([0, VIDEO_H - corner, 2, VIDEO_H],       fill=(*color, ca))
    draw.rectangle([VIDEO_W - corner, VIDEO_H - 2, VIDEO_W, VIDEO_H], fill=(*color, ca))
    draw.rectangle([VIDEO_W - 2, VIDEO_H - corner, VIDEO_W, VIDEO_H], fill=(*color, ca))

    return np.array(frame.convert("RGB"))


# ═════════════════════════════════════════════════════════════════════════════
# PPTX → GÖRÜNTÜLER
# ═════════════════════════════════════════════════════════════════════════════
def pptx_to_images(pptx_bytes: bytes) -> list:
    tmp = tempfile.mkdtemp(prefix="pptx2img_")
    try:
        pptx_path = os.path.join(tmp, "presentation.pptx")
        with open(pptx_path, "wb") as f:
            f.write(pptx_bytes)
        try:
            _run([LO_BIN, "--headless", "--convert-to", "pdf",
                  "--outdir", tmp, pptx_path],
                 timeout=180, step_name="LibreOffice PDF")
        except RuntimeError as e:
            raise RuntimeError(f"LibreOffice başarısız.\npackages.txt: libreoffice\n\n{e}")
        pdfs = [f for f in os.listdir(tmp) if f.endswith(".pdf")]
        if not pdfs:
            raise RuntimeError("LibreOffice PDF üretmedi.")
        pdf_path = os.path.join(tmp, pdfs[0])
        img_prefix = os.path.join(tmp, "slide")
        try:
            _run(["pdftoppm", "-jpeg", "-r", "192", pdf_path, img_prefix],
                 timeout=120, step_name="pdftoppm")
        except RuntimeError as e:
            raise RuntimeError(f"pdftoppm başarısız.\npackages.txt: poppler-utils\n\n{e}")
        files = sorted([
            os.path.join(tmp, f) for f in os.listdir(tmp)
            if f.startswith("slide") and (f.endswith(".jpg") or f.endswith(".jpeg"))
        ])
        if not files:
            raise RuntimeError("pdftoppm görüntü üretmedi.")
        images = []
        for p in files:
            src = Image.open(p).convert("RGB")
            sw, sh = src.size
            canvas = Image.new("RGB", (VIDEO_W, VIDEO_H), (6, 10, 8))
            scale  = min(SLIDE_AREA_W / sw, SLIDE_AREA_H / sh)
            nw, nh = int(sw * scale), int(sh * scale)
            paste_x = (SLIDE_AREA_W - nw) // 2
            paste_y = TOP_BAR + (SLIDE_AREA_H - nh) // 2
            canvas.paste(src.resize((nw, nh), Image.LANCZOS), (paste_x, paste_y))
            images.append(canvas)
        return images
    finally:
        shutil.rmtree(tmp, ignore_errors=True)

def read_pptx_notes(pptx_bytes: bytes) -> list:
    prs = Presentation(io.BytesIO(pptx_bytes))
    notes = []
    for slide in prs.slides:
        txt = ""
        try:
            txt = slide.notes_slide.notes_text_frame.text.strip()
        except:
            pass
        notes.append(txt)
    return notes

# ═════════════════════════════════════════════════════════════════════════════
# SES HAZIRLAMA
# ═════════════════════════════════════════════════════════════════════════════
def prepare_audio_segments(
    slide_audio_map, durations, n_slides, global_audio, use_global, work_dir
):
    audio_paths, seek_starts, dur_list = [], [], []

    if use_global and global_audio:
        raw_g = os.path.join(work_dir, "global_raw.audio")
        with open(raw_g, "wb") as f:
            f.write(global_audio)
        raw_dur = audio_duration_ffprobe(raw_g)
        clean_g = os.path.join(work_dir, "global_clean.aac")
        clean_audio(raw_g, clean_g, step="Global ses temizleme")
        clean_dur = audio_duration_ffprobe(clean_g)
        total_audio_dur = min(raw_dur, clean_dur)

        user_durs  = [durations.get(i, 0.0) for i in range(n_slides)]
        total_user = sum(user_durs)
        if total_user > 0.5:
            seg_durs = [total_audio_dur * (d / total_user) for d in user_durs]
        else:
            per = total_audio_dur / max(n_slides, 1)
            seg_durs = [per] * n_slides

        acc = 0.0
        for sd in seg_durs:
            audio_paths.append(clean_g)
            seek_starts.append(acc)
            dur_list.append(sd)
            acc += sd
    else:
        for i in range(n_slides):
            ab = slide_audio_map.get(i)
            if ab:
                raw   = os.path.join(work_dir, f"raw_{i:04d}.audio")
                clean = os.path.join(work_dir, f"seg_{i:04d}.aac")
                with open(raw, "wb") as f:
                    f.write(ab)
                try:
                    clean_audio(raw, clean, step=f"Slayt {i+1} ses")
                    real_dur = audio_duration_ffprobe(clean)
                    audio_paths.append(clean)
                    seek_starts.append(0.0)
                    dur_list.append(real_dur)
                except:
                    audio_paths.append(None)
                    seek_starts.append(0.0)
                    dur_list.append(durations.get(i, 3.0))
            else:
                audio_paths.append(None)
                seek_starts.append(0.0)
                dur_list.append(durations.get(i, 3.0))

    return audio_paths, seek_starts, dur_list

# ═════════════════════════════════════════════════════════════════════════════
# KARE RENDER
# ═════════════════════════════════════════════════════════════════════════════
def _draw_speaker_dot(draw, x, y, r, color):
    draw.ellipse([x-r, y-r, x+r, y+r], fill=color)

def render_frame(slide_img, slide_idx, total, t, speaker, has_audio):
    frame = slide_img.copy()
    draw  = ImageDraw.Draw(frame, "RGBA")
    wi, hi = frame.size
    color = speaker.get("rgb", BRAND_RGB)
    name  = speaker.get("name", "Elif Aracıoğlu")
    role  = speaker.get("role", "Eczacı")

    fn18 = _font(18)
    fn13 = _font(13)
    fn11 = _font(11)

    # Üst bant
    draw.rectangle([0, 0, wi, TOP_BAR], fill=(4, 8, 6, 255))
    draw.rectangle([0, TOP_BAR-2, wi, TOP_BAR], fill=(*color, 255))
    draw.text((16, 8),  "POL",   font=fn18, fill=(*color, 255))
    draw.text((16, 32), "CAST",  font=fn11, fill=(*color, 170))
    draw.rectangle([108, 10, 110, TOP_BAR-8], fill=(*color, 50))
    dot_cx, dot_cy = 122, TOP_BAR // 2
    _draw_speaker_dot(draw, dot_cx, dot_cy, 6, color)
    draw.text((dot_cx+12, dot_cy-7), f"{name}  |  {role}",
              font=fn13, fill=(195, 235, 215, 220))
    da = int(175 + 80 * math.sin(t * math.pi * 4))
    cx = wi // 2 + 80
    draw.ellipse([cx-6, dot_cy-6, cx+6, dot_cy+6], fill=(210, 55, 55, da))
    draw.text((cx+12, dot_cy-7), "CANLI", font=fn11, fill=(210, 55, 55, 210))
    brand_text = "Eczacı Elif Aracıoglu"
    try:
        btw = draw.textlength(brand_text, font=fn11)
    except:
        btw = len(brand_text) * 6
    bx = wi - int(btw) - 28
    _draw_speaker_dot(draw, bx-8, dot_cy, 5, color)
    draw.text((bx, dot_cy-7), brand_text, font=fn11, fill=(*color, 145))

    # Alt bant
    bot_y = hi - BOT_BAR
    draw.rectangle([0, bot_y, wi, hi], fill=(4, 8, 6, 255))
    draw.rectangle([0, bot_y, wi, bot_y+2], fill=(*color, 255))
    draw.text((16, bot_y+(BOT_BAR-14)//2), f"Slayt {slide_idx+1}  /  {total}",
              font=fn13, fill=(140, 205, 175, 210))
    mid_text = "POLCAST"
    try:
        mw = draw.textlength(mid_text, font=fn11)
    except:
        mw = len(mid_text) * 6
    draw.text((wi//2-int(mw)//2, bot_y+(BOT_BAR-12)//2),
              mid_text, font=fn11, fill=(*color, 90))
    pw = int(wi * (slide_idx + t) / max(total, 1))
    draw.rectangle([0, hi-5, wi, hi], fill=(8, 14, 11, 255))
    draw.rectangle([0, hi-5, pw, hi], fill=(*color, 255))
    if has_audio:
        bc, bw2, bg = 9, 4, 4
        bx0 = wi - bc*(bw2+bg) - 16
        by  = hi - 8
        for bi in range(bc):
            bh  = int(3 + 14*abs(math.sin(t*math.pi*4.2+bi*0.95)))
            bx2 = bx0 + bi*(bw2+bg)
            draw.rounded_rectangle([bx2, by-bh, bx2+bw2, by], radius=2, fill=(*color, 255))

    return np.array(frame.convert("RGB"))

# ═════════════════════════════════════════════════════════════════════════════
# VIDEO OLUŞTURMA
# ═════════════════════════════════════════════════════════════════════════════
def _pipe_frames_to_file(frames_iter, nf: int, out_path: str):
    cmd = [
        FFMPEG, "-y",
        "-f", "rawvideo", "-vcodec", "rawvideo",
        "-s", f"{VIDEO_W}x{VIDEO_H}", "-pix_fmt", "rgb24",
        "-r", str(VIDEO_FPS), "-i", "pipe:0",
        "-vcodec", "libx264", "-crf", "22", "-preset", "fast",
        "-pix_fmt", "yuv420p", "-r", str(VIDEO_FPS), "-vsync", "cfr",
        "-an", out_path,
    ]
    try:
        proc = subprocess.Popen(cmd, stdin=subprocess.PIPE,
                                stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
    except FileNotFoundError:
        raise RuntimeError(f"ffmpeg bulunamadı: {FFMPEG}")
    try:
        for frame in frames_iter:
            proc.stdin.write(frame.astype(np.uint8).tobytes())
        proc.stdin.close()
        proc.wait(timeout=600)
        if proc.returncode != 0:
            raise RuntimeError(f"Pipe encode başarısız: {out_path}")
    except Exception as e:
        proc.kill()
        raise RuntimeError(f"Pipe encode hatası: {e}")

def _mux(video_path: str, audio_path: str, out_path: str):
    _run([
        FFMPEG, "-y",
        "-i", video_path,
        "-i", audio_path,
        "-map", "0:v:0", "-map", "1:a:0",
        "-c:v", "copy",
        "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2",
        "-shortest",
        "-movflags", "+faststart",
        out_path,
    ], timeout=300, step_name="Mux")

def _silence(work_dir: str, dur: float, tag: str) -> str:
    out = os.path.join(work_dir, f"silence_{tag}.aac")
    _run([
        FFMPEG, "-y", "-f", "lavfi",
        "-i", "anullsrc=r=44100:cl=stereo",
        "-t", str(dur),
        "-c:a", "aac", "-b:a", "128k", "-ar", "44100", "-ac", "2",
        out,
    ], timeout=30, step_name=f"Sessizlik {tag}")
    return out

def _concat_videos(parts: list, out_path: str, work_dir: str):
    concat_list = os.path.join(work_dir, "concat.txt")
    with open(concat_list, "w") as f:
        for p in parts:
            f.write(f"file '{p}'\n")
    _run([
        FFMPEG, "-y", "-f", "concat", "-safe", "0", "-i", concat_list,
        "-c", "copy", "-movflags", "+faststart", out_path,
    ], timeout=600, step_name="Concat")

def build_video(slide_images, audio_paths, seek_starts, durations, speakers, work_dir, cb=None):
    if not FFMPEG or not os.path.exists(FFMPEG):
        raise RuntimeError("ffmpeg bulunamadı!")

    n          = len(slide_images)
    audio_file = next((p for p in audio_paths if p and os.path.exists(p)), None)
    has_audio  = audio_file is not None
    color      = speakers[0].get("rgb", BRAND_RGB) if speakers else BRAND_RGB

    jingle_nf  = max(1, round(JINGLE_DUR * VIDEO_FPS))

    # 1. Açılış jingle video
    if cb: cb(0.03, "Açılış jingle hazırlanıyor…")
    intro_vid = os.path.join(work_dir, "intro_vid.mp4")
    def intro_frames():
        for fi in range(jingle_nf):
            t = fi / max(jingle_nf - 1, 1)
            yield render_intro_frame(t, "open", color)
    _pipe_frames_to_file(intro_frames(), jingle_nf, intro_vid)

    # 2. Açılış ses
    if cb: cb(0.08, "Açılış jingle sesi sentezleniyor…")
    intro_audio = make_jingle(work_dir, "open")
    intro_muxed = os.path.join(work_dir, "intro_muxed.mp4")
    _mux(intro_vid, intro_audio, intro_muxed)

    # 3. Ana içerik video
    if cb: cb(0.15, "Ana içerik encode ediliyor…")
    nf_list  = [max(1, round(d * VIDEO_FPS)) for d in durations]
    main_vid = os.path.join(work_dir, "main_vid.mp4")
    def main_frames():
        for slide_idx, (img, spk, nf) in enumerate(zip(slide_images, speakers, nf_list)):
            if cb:
                cb(0.15 + 0.60 * (slide_idx / n),
                   f"Slayt {slide_idx+1}/{n} — {nf/VIDEO_FPS:.1f}sn")
            for fi in range(nf):
                t = fi / max(nf - 1, 1)
                yield render_frame(img, slide_idx, n, t, spk, has_audio)
    _pipe_frames_to_file(main_frames(), sum(nf_list), main_vid)

    # 4. Ana ses
    if cb: cb(0.77, "Ana ses birleştiriliyor…")
    main_audio = audio_file if has_audio else _silence(work_dir, sum(durations) or n * 3.0, "main")
    main_muxed = os.path.join(work_dir, "main_muxed.mp4")
    _mux(main_vid, main_audio, main_muxed)

    # 5. Kapanış jingle video
    if cb: cb(0.82, "Kapanış jingle hazırlanıyor…")
    outro_vid = os.path.join(work_dir, "outro_vid.mp4")
    def outro_frames():
        for fi in range(jingle_nf):
            t = fi / max(jingle_nf - 1, 1)
            yield render_intro_frame(t, "close", color)
    _pipe_frames_to_file(outro_frames(), jingle_nf, outro_vid)

    # 6. Kapanış ses
    if cb: cb(0.87, "Kapanış jingle sesi sentezleniyor…")
    outro_audio = make_jingle(work_dir, "close")
    outro_muxed = os.path.join(work_dir, "outro_muxed.mp4")
    _mux(outro_vid, outro_audio, outro_muxed)

    # 7. Birleştir: Açılış + Ana + Kapanış
    if cb: cb(0.93, "Açılış + içerik + kapanış birleştiriliyor…")
    out_mp4 = os.path.join(work_dir, "output.mp4")
    _concat_videos([intro_muxed, main_muxed, outro_muxed], out_mp4, work_dir)

    for tmp in [intro_vid, intro_audio, intro_muxed,
                main_vid, main_muxed,
                outro_vid, outro_audio, outro_muxed]:
        try: os.unlink(tmp)
        except: pass

    if cb: cb(1.0, "Tamamlandı! ✅")
    if os.path.exists(out_mp4):
        with open(out_mp4, "rb") as f:
            return f.read()
    raise RuntimeError("Çıktı MP4 oluşturulamadı.")


CSS = """
<style>
@import url('https://fonts.googleapis.com/css2?family=Cormorant+Garamond:ital,wght@0,400;0,600;0,700;1,400;1,600&family=DM+Sans:opsz,wght@9..40,300;9..40,400;9..40,500;9..40,600&display=swap');
:root{
  --g:#34A883;--gd:rgba(52,168,131,.15);--gg:rgba(52,168,131,.055);
  --gl:rgba(52,168,131,.32);
  --bg:#060d0a;--sf:rgba(255,255,255,.024);--br:rgba(255,255,255,.065);
  --tx:#ddeee8;--mu:#42665a;--mu2:#2a4038;
}
html,body,[class*="css"]{font-family:'DM Sans',sans-serif;background:var(--bg);}
.stApp{
  background:
    radial-gradient(ellipse 75% 50% at 5%  5%, rgba(52,168,131,.06) 0%,transparent 55%),
    radial-gradient(ellipse 55% 45% at 95% 95%,rgba(52,168,131,.04) 0%,transparent 55%),
    var(--bg);
  color:var(--tx);
}
section[data-testid="stSidebar"]{
  background:rgba(4,9,7,.98);border-right:1px solid var(--gd);
}
.hero{
  position:relative;overflow:hidden;padding:3rem 1rem 1.8rem;
  border-bottom:1px solid var(--gd);margin-bottom:1.8rem;text-align:center;
}
.hero::before{
  content:'';position:absolute;inset:0;pointer-events:none;
  background:radial-gradient(ellipse 65% 100% at 50% -5%,rgba(52,168,131,.07),transparent 60%);
}
.hero-pill{
  display:inline-flex;align-items:center;gap:.4rem;padding:.24rem .85rem;
  border-radius:50px;background:var(--gg);border:1px solid var(--gl);
  font-size:.62rem;letter-spacing:.22em;text-transform:uppercase;
  color:var(--g);margin-bottom:1.1rem;font-weight:600;
}
.hero h1{
  font-family:'Cormorant Garamond',serif;
  font-size:3.2rem;font-weight:700;line-height:1.05;
  color:#e8f5ef;margin:0 0 .15rem;letter-spacing:-.02em;
}
.hero h1 em{font-style:italic;font-weight:400;color:var(--g);}
.hero-author{
  font-size:.72rem;letter-spacing:.22em;text-transform:uppercase;
  color:var(--mu);margin-top:.65rem;
}
.hero-author strong{color:var(--g);font-weight:600;letter-spacing:.18em;}
.hero-rule{
  width:52px;height:1px;margin:1rem auto 0;
  background:linear-gradient(90deg,transparent,var(--g),transparent);
}
.lbl{
  font-size:.58rem;letter-spacing:.24em;text-transform:uppercase;
  color:var(--mu);margin:.9rem 0 .45rem;
  padding-left:.55rem;border-left:2px solid var(--g);
}
.srow{
  display:flex;gap:.9rem;flex-wrap:wrap;padding:.55rem 1rem;margin:.5rem 0;
  background:var(--gg);border:1px solid var(--gd);border-radius:10px;
  font-size:.78rem;color:var(--mu);align-items:center;
}
.srow strong{color:var(--g);}
.char-card{
  display:flex;align-items:center;gap:.8rem;padding:.65rem 1rem;margin:.32rem 0;
  background:var(--sf);border:1px solid var(--br);border-radius:11px;
}
.char-dot{width:11px;height:11px;border-radius:50%;flex-shrink:0;}
.char-name{font-size:.88rem;font-weight:600;color:#c8e8d8;}
.char-role{font-size:.7rem;color:var(--mu);margin-left:auto;}
.sl-card{
  padding:.6rem .9rem;margin:.26rem 0;background:var(--sf);
  border:1px solid var(--br);border-radius:9px;border-left:3px solid;
}
.sl-num{font-size:.63rem;font-weight:700;letter-spacing:.14em;text-transform:uppercase;margin-bottom:.18rem;}
.sl-note{font-size:.78rem;color:#5a8070;line-height:1.5;font-style:italic;}
.badge{display:inline-flex;align-items:center;gap:.28rem;
  padding:.18rem .55rem;border-radius:50px;font-size:.65rem;font-weight:700;}
.badge-gold{background:rgba(52,168,131,.1);color:var(--g);border:1px solid rgba(52,168,131,.22);}
.dep{font-size:.72rem;margin:.1rem 0;line-height:1.55;}
.ok{color:#55c98a;}.er{color:#d97070;}
.ffmpeg-warn{
  padding:.8rem 1rem;margin:.6rem 0;border-radius:10px;
  background:rgba(217,112,112,.07);border:1px solid rgba(217,112,112,.22);
  font-size:.78rem;color:#d97070;line-height:1.6;
}
.ffmpeg-warn code{background:rgba(255,255,255,.07);padding:.1rem .3rem;border-radius:4px;font-size:.73rem;}
.audio-info{
  padding:.5rem .85rem;margin:.4rem 0;
  background:rgba(52,168,131,.05);border:1px solid rgba(52,168,131,.14);
  border-radius:8px;font-size:.72rem;color:var(--mu);line-height:1.6;
}
.audio-info strong{color:var(--g);}
.slide-info{
  padding:.6rem 1rem;margin:.5rem 0;border-radius:8px;
  background:rgba(52,168,131,.08);border:1px solid rgba(52,168,131,.25);
  font-size:.75rem;color:#7cc8a8;line-height:1.6;
}
.slide-info strong{color:var(--g);}
.jingle-info{
  padding:.6rem 1rem;margin:.5rem 0;border-radius:8px;
  background:rgba(52,168,131,.06);border:1px solid rgba(52,168,131,.20);
  font-size:.74rem;color:#6ab898;line-height:1.75;
}
.jingle-info strong{color:var(--g);}
.dl-alt{
  display:block;text-align:center;padding:9px 14px;margin-top:8px;
  background:rgba(52,168,131,.07);border:1px solid rgba(52,168,131,.2);
  color:var(--g);border-radius:9px;font-weight:600;font-size:.8rem;
  text-decoration:none;
}
.stButton>button[kind="primary"]{
  background:linear-gradient(135deg,#289970,#3ec49a)!important;
  color:#030a07!important;font-weight:700!important;border:none!important;letter-spacing:.04em;
}
.stButton>button[kind="primary"]:hover{
  filter:brightness(1.07);transform:translateY(-1px);
  box-shadow:0 4px 22px rgba(52,168,131,.3)!important;
}
.stProgress>div>div{border-radius:10px;}
.sb-brand{
  text-align:center;padding:1.15rem 0 .9rem;
  border-bottom:1px solid var(--gd);margin-bottom:1rem;
}
.sb-title{
  font-family:'Cormorant Garamond',serif;
  font-size:1.25rem;font-weight:700;color:#e0f0e8;letter-spacing:-.01em;
}
.sb-title em{font-style:italic;font-weight:400;color:var(--g);}
.sb-name{font-size:.61rem;color:var(--mu);letter-spacing:.2em;text-transform:uppercase;margin-top:.35rem;}
.sb-name strong{color:var(--g);font-weight:600;}
audio{width:100%;border-radius:8px;margin:3px 0;}
hr{border-color:var(--br);}
input[type="text"]{
  background:rgba(255,255,255,.04)!important;
  border:1px solid rgba(255,255,255,.08)!important;
  color:var(--tx)!important;border-radius:8px!important;
}
</style>
"""

# ═════════════════════════════════════════════════════════════════════════════
# SESSION STATE
# ═════════════════════════════════════════════════════════════════════════════
def init_state():
    defaults = {
        "ss_pptx_bytes":    None,
        "ss_slide_notes":   [],
        "ss_slide_images":  [],
        "ss_slide_audio":   {},
        "ss_slide_speaker": {},
        "ss_global_audio":  None,
        "ss_durations":     {},
        "ss_use_global":    True,
        "ss_video_bytes":   None,
        "ss_characters":    [c.copy() for c in DEFAULT_CHARACTERS],
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v

# ═════════════════════════════════════════════════════════════════════════════
# KARAKTERLERİ YÖNETME
# ═════════════════════════════════════════════════════════════════════════════
def character_manager_ui():
    chars = st.session_state.ss_characters
    st.markdown('<p class="lbl">🎭 Konuşmacılar</p>', unsafe_allow_html=True)
    st.caption("Elif Aracıoğlu varsayılan konuşmacıdır. İstediğiniz kadar ekleyip çıkarabilirsiniz.")
    to_delete = None
    for i, ch in enumerate(chars):
        is_default = (i == 0)
        with st.container():
            c1, c2, c3, c4 = st.columns([2.5, 2, 1.2, 0.7])
            with c1:
                new_name = st.text_input(
                    "İsim", value=ch["name"], key=f"ch_name_{i}",
                    label_visibility="collapsed", disabled=is_default, placeholder="İsim...")
                if not is_default:
                    chars[i]["name"] = new_name or ch["name"]
            with c2:
                new_role = st.text_input(
                    "Rol", value=ch["role"], key=f"ch_role_{i}",
                    label_visibility="collapsed", placeholder="Rol...")
                chars[i]["role"] = new_role or ch["role"]
            with c3:
                emojis = ["💊","🎤","👩‍💼","🎧","🎙️","💬","📢","🗣️","👤","🎵"]
                cur_e  = ch.get("emoji", "💊")
                sel_e  = st.selectbox(
                    "Emoji", emojis,
                    index=emojis.index(cur_e) if cur_e in emojis else 0,
                    key=f"ch_emoji_{i}", label_visibility="collapsed")
                chars[i]["emoji"] = sel_e
            with c4:
                if is_default:
                    st.markdown('<span class="badge badge-gold">Varsayılan</span>',
                                unsafe_allow_html=True)
                else:
                    st.markdown("<br>", unsafe_allow_html=True)
                    if st.button("✕", key=f"ch_del_{i}",
                                 help=f"{ch['name']} sil", use_container_width=True):
                        to_delete = i
            st.markdown(
                f'<div class="char-card" style="margin-top:-6px;padding:.4rem .9rem;">'
                f'<div class="char-dot" style="background:#{ch["hex"]};"></div>'
                f'<span class="char-name">{ch["emoji"]}  {chars[i]["name"]}</span>'
                f'<span class="char-role">{chars[i]["role"]}</span></div>',
                unsafe_allow_html=True)
    if to_delete is not None:
        chars.pop(to_delete)
        new_map = {}
        for k, v in st.session_state.ss_slide_speaker.items():
            new_map[k] = v if v < to_delete else max(0, v - 1)
        st.session_state.ss_slide_speaker = new_map
        st.rerun()
    st.markdown("---")
    with st.expander("➕ Yeni Konuşmacı Ekle", expanded=False):
        c_n, c_r, c_add = st.columns([2, 2, 1])
        with c_n:
            nn = st.text_input("İsim", key="wu_new_name", placeholder="Örn: Ecem")
        with c_r:
            nr = st.text_input("Rol",  key="wu_new_role", placeholder="Örn: Konuk")
        with c_add:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Ekle ➕", use_container_width=True, key="wu_add_char"):
                if nn.strip():
                    chars.append({**PALETTE[len(chars) % len(PALETTE)],
                                  "name": nn.strip(), "role": nr.strip() or "Konuşmacı"})
                    st.rerun()
                else:
                    st.warning("İsim boş olamaz.")
    st.session_state.ss_characters = chars
    return chars

# ═════════════════════════════════════════════════════════════════════════════
# SES ATAMA
# ═════════════════════════════════════════════════════════════════════════════
def audio_assignment_ui(n_slides: int, chars: list):
    mode = st.radio(
        "Ses Modu",
        ["🔊 Tek ses — tüm sunuma", "🎙️ Her slayta ayrı ses"],
        key="wu_mode_radio", horizontal=True, label_visibility="collapsed",
    )
    use_global = mode.startswith("🔊")
    st.session_state.ss_use_global = use_global

    st.markdown(
        '<div class="audio-info">'
        '🧹 <strong>Otomatik ses temizleme:</strong> '
        'Alçak frekans gürültüsü (highpass 80 Hz), arka plan sesi (afftdn −20 dB) '
        've ses seviyesi dengeleme uygulanır. '
        'Seste atlama olmaz — temiz, kesintisiz ses.'
        '</div>',
        unsafe_allow_html=True,
    )

    if use_global:
        st.caption(
            "**Tek ses modu:** Ses önce temizlenir, ardından slayt sayısına eşit bölünür."
        )
        col_g1, col_g2 = st.columns([2, 1])
        with col_g1:
            gf = st.file_uploader(
                "Genel ses", type=["mp3","wav","m4a","ogg"],
                key="wu_glob_upload", label_visibility="collapsed",
            )
            if gf is not None:
                ab = gf.read()
                st.session_state.ss_global_audio = ab
                total_dur = audio_duration_sec_bytes(ab)
                per_slide = total_dur / max(n_slides, 1)
                for i in range(n_slides):
                    if st.session_state.ss_durations.get(i, 0.0) <= 0.5:
                        st.session_state.ss_durations[i] = per_slide
                st.audio(ab, format="audio/mp3")
                cur_weights = [st.session_state.ss_durations.get(i, per_slide) for i in range(n_slides)]
                cur_total   = sum(cur_weights)
                st.caption(
                    f"Toplam ses ~{total_dur:.1f} sn  ·  "
                    f"Slayt süre toplamı: {cur_total:.1f} sn"
                )
        with col_g2:
            st.markdown('<p class="lbl">Tüm Slaytlar — Konuşmacı</p>', unsafe_allow_html=True)
            char_names = [f'{c["emoji"]} {c["name"]}' for c in chars]
            sel = st.selectbox(
                "Konuşmacı", char_names, index=0,
                key="wu_global_speaker", label_visibility="collapsed",
            )
            sel_idx = char_names.index(sel)
            for i in range(n_slides):
                st.session_state.ss_slide_speaker[i] = sel_idx
    else:
        st.caption(
            "**Slayt bazlı mod:** Her slayta farklı ses ve konuşmacı atayabilirsiniz."
        )
        char_names = [f'{c["emoji"]} {c["name"]}' for c in chars]
        notes = st.session_state.ss_slide_notes
        cqa1, cqa2 = st.columns([3, 1])
        with cqa1:
            qs = st.selectbox("Tüm slaytlara ata:", char_names, index=0, key="wu_quick_speaker")
        with cqa2:
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("Tümüne Uygula", key="wu_quick_apply", use_container_width=True):
                qi = char_names.index(qs)
                for i in range(n_slides):
                    st.session_state.ss_slide_speaker[i] = qi
                st.rerun()
        st.markdown("---")
        for i in range(n_slides):
            note_preview = (notes[i][:80]+"…") if i < len(notes) and len(notes[i])>80 \
                else (notes[i] if i < len(notes) else "")
            ci = min(st.session_state.ss_slide_speaker.get(i, 0), len(chars)-1)
            sc = f'#{chars[ci]["hex"]}'
            st.markdown(
                f'<div class="sl-card" style="border-left-color:{sc};">'
                f'<div class="sl-num" style="color:{sc};">'
                f'Slayt {i+1}  ·  {chars[ci]["emoji"]} {chars[ci]["name"]}</div>'
                f'<div class="sl-note">{note_preview or "—"}</div></div>',
                unsafe_allow_html=True)
            sc1, sc2, sc3 = st.columns([1.5, 2.5, 0.5])
            with sc1:
                sel = st.selectbox(f"Konuşmacı S{i+1}", char_names,
                                   index=ci, key=f"wu_spk_{i}", label_visibility="collapsed")
                st.session_state.ss_slide_speaker[i] = char_names.index(sel)
            with sc2:
                uf = st.file_uploader(f"Ses S{i+1}", type=["mp3","wav","m4a","ogg"],
                                      key=f"wu_sl_{i}", label_visibility="collapsed")
                if uf is not None:
                    ab = uf.read()
                    st.session_state.ss_slide_audio[i] = ab
                    dur = audio_duration_sec_bytes(ab)
                    st.session_state.ss_durations[i] = dur
                    st.audio(ab, format="audio/mp3")
            with sc3:
                dv   = st.session_state.ss_durations.get(i, 3.0)
                icon = "🔊" if i in st.session_state.ss_slide_audio else "🔇"
                st.caption(f"{icon} {dv:.1f}s")

    with st.expander("⚙️ Slayt sürelerini manuel ayarla (opsiyonel)", expanded=False):
        per_row   = min(n_slides, 5)
        dur_grids = [st.columns(per_row) for _ in range(-(-n_slides // per_row))]
        for i in range(n_slides):
            r, c = i // per_row, i % per_row
            with dur_grids[r][c]:
                d = st.number_input(
                    f"S{i+1} (sn)", min_value=0.5, max_value=300.0,
                    value=float(st.session_state.ss_durations.get(i, 3.0)),
                    step=0.5, key=f"wu_dur_{i}")
                st.session_state.ss_durations[i] = d
    return use_global

# ═════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ═════════════════════════════════════════════════════════════════════════════
def render_sidebar():
    with st.sidebar:
        st.markdown(
            '<div class="sb-brand">'
            '<div class="sb-title">3 <em>Soru</em> 3 Dakika</div>'
            '<div class="sb-name">Eczacı <strong>Elif Aracıoğlu</strong></div>'
            '</div>',
            unsafe_allow_html=True,
        )
        st.markdown("---")
        character_manager_ui()
        st.markdown("---")
        st.markdown('<p class="lbl">Sistem Durumu</p>', unsafe_allow_html=True)
        ffmpeg_hint  = "requirements.txt: imageio[ffmpeg]  VEYA  packages.txt: ffmpeg"
        ffmpeg_label = "ffmpeg"
        if FFMPEG_OK:
            ffmpeg_label = (f"ffmpeg <span style='color:#2a4038;font-size:.6rem;'>"
                            f"({FFMPEG})</span>")
        checks = [
            ("Pillow",       PIL_OK,    "pip: pillow"),
            ("imageio",      IMAGEIO_OK,"pip: imageio[ffmpeg]"),
            ("python-pptx",  PPTX_OK,   "pip: python-pptx"),
            ("LibreOffice",  LO_OK,     "packages.txt: libreoffice"),
            ("pdftoppm",     PPM_OK,    "packages.txt: poppler-utils"),
            ("ffmpeg",       FFMPEG_OK, ffmpeg_hint),
        ]
        for name, ok, hint in checks:
            cls  = "ok" if ok else "er"
            icon = "🟢" if ok else "🔴"
            extra = f' <span style="color:#2a4038;font-size:.62rem;">— {hint}</span>' if not ok else ""
            dn = ffmpeg_label if name == "ffmpeg" else name
            st.markdown(f'<div class="dep {cls}">{icon} {dn}{extra}</div>',
                        unsafe_allow_html=True)
        if not FFMPEG_OK:
            st.markdown(
                '<div class="ffmpeg-warn">⚠️ <b>ffmpeg bulunamadı!</b><br>'
                'requirements.txt:<br><code>imageio[ffmpeg]</code><br><br>'
                'VEYA packages.txt:<br><code>ffmpeg</code></div>',
                unsafe_allow_html=True)
        if not LO_OK or not PPM_OK:
            st.markdown("---")
            st.caption("**packages.txt**:\n```\nlibreoffice\npoppler-utils\n```")
        st.markdown("---")
        chars = st.session_state.ss_characters
        st.markdown('<p class="lbl">Aktif Konuşmacılar</p>', unsafe_allow_html=True)
        for ch in chars:
            st.markdown(
                f'<div style="display:flex;align-items:center;gap:.5rem;'
                f'padding:.2rem 0;font-size:.78rem;">'
                f'<div style="width:9px;height:9px;border-radius:50%;'
                f'background:#{ch["hex"]};"></div>'
                f'<span style="color:#c8e8d8;">{ch["emoji"]} {ch["name"]}</span>'
                f'<span style="color:#2a4038;margin-left:auto;font-size:.65rem;">'
                f'{ch["role"]}</span></div>',
                unsafe_allow_html=True)
        st.markdown("---")
        st.markdown(
            f'<div style="font-size:.6rem;color:#2a4038;line-height:1.7;">'
            f'📐 Video: {VIDEO_W}×{VIDEO_H} · {VIDEO_FPS}fps<br>'
            f'🟩 Slayt alanı: {SLIDE_AREA_W}×{SLIDE_AREA_H}px<br>'
            f'📌 Üst bant: {TOP_BAR}px · Alt bant: {BOT_BAR}px<br>'
            f'🎵 Jingle: {int(JINGLE_DUR)}sn açılış + {int(JINGLE_DUR)}sn kapanış<br>'
            f'✅ Başlıklar tam görünür — bant dışı'
            f'</div>',
            unsafe_allow_html=True)
        st.markdown("---")
        st.markdown(
            '<div style="font-size:.6rem;color:#2a4038;text-align:center;">'
            'v16.0 · Medikal Jingle · EKG Animasyon · Temiz Ses</div>',
            unsafe_allow_html=True)

# ═════════════════════════════════════════════════════════════════════════════
# MAIN
# ═════════════════════════════════════════════════════════════════════════════
def main():
    st.set_page_config(
        page_title="3 Soru 3 Dakika · Elif Aracıoğlu",
        page_icon="💊", layout="wide",
        initial_sidebar_state="expanded",
    )
    st.markdown(CSS, unsafe_allow_html=True)
    init_state()
    render_sidebar()

    st.markdown(
        '<div class="hero">'
        '<div class="hero-pill">💊 Video Stüdyo</div>'
        '<h1>3 <em>Soru</em><br>3 Dakika</h1>'
        '<div class="hero-author">Eczacı &nbsp;<strong>Elif Aracıoğlu</strong></div>'
        '<div class="hero-rule"></div>'
        '</div>',
        unsafe_allow_html=True,
    )

    chars = st.session_state.ss_characters

    # ── ADIM 1: PPTX ──────────────────────────────────────────────────────────
    st.markdown('<p class="lbl">① PowerPoint Dosyası</p>', unsafe_allow_html=True)

    st.markdown(
        f'<div class="slide-info">'
        f'✅ <strong>Slayt içeriği bantların dışında:</strong> '
        f'Üst bant ({TOP_BAR}px) ve alt bant ({BOT_BAR}px) slayt alanının dışındadır. '
        f'Slayt görüntüsü {SLIDE_AREA_W}×{SLIDE_AREA_H}px alana tam sığdırılır.'
        f'</div>',
        unsafe_allow_html=True,
    )

    # Jingle bilgi kutusu
    st.markdown(
        '<div class="jingle-info">'
        '🎵 <strong>Profesyonel Medikal Jingle v2.0</strong> — Her video otomatik olarak:<br>'
        '&nbsp;&nbsp;▶ <strong>Açılış (5sn):</strong> EKG bip zinciri · Steril sine pad · Diapason 440 Hz · Kristal arpej C5→C6<br>'
        '&nbsp;&nbsp;◀ <strong>Kapanış (5sn):</strong> Yavaşlayan EKG · İnen arpej C6→C5 · Medikal onay çifti 1320+1760 Hz<br>'
        '&nbsp;&nbsp;✨ Tüm sesler sıfırdan sentezlenir — dış dosya gerekmez'
        '</div>',
        unsafe_allow_html=True,
    )

    col_a, col_b = st.columns([1, 1], gap="large")
    with col_a:
        pptx_file = st.file_uploader(
            "PPTX", type=["pptx"], key="wu_pptx", label_visibility="collapsed")
        if pptx_file is not None:
            raw = pptx_file.read()
            if raw != st.session_state.ss_pptx_bytes:
                st.session_state.ss_pptx_bytes   = raw
                st.session_state.ss_slide_images  = []
                st.session_state.ss_video_bytes   = None
                st.session_state.ss_slide_audio   = {}
                st.session_state.ss_slide_speaker = {}
                st.session_state.ss_durations     = {}
                with st.spinner("PPTX okunuyor..."):
                    try:
                        st.session_state.ss_slide_notes = read_pptx_notes(raw)
                    except Exception as e:
                        st.error(f"PPTX okunamadı: {e}")
                        st.session_state.ss_slide_notes = []
        if st.session_state.ss_pptx_bytes:
            n = len(st.session_state.ss_slide_notes)
            st.success(f"✅ Yüklendi — **{n}** slayt")
            char_names = [f'{c["emoji"]} {c["name"]}' for c in chars]
            spk_counts = {}
            for i in range(n):
                idx = min(st.session_state.ss_slide_speaker.get(i, 0), len(chars)-1)
                nm  = chars[idx]["name"]
                spk_counts[nm] = spk_counts.get(nm, 0) + 1
            badges = " ".join(
                f'<span class="badge badge-gold">{nm}: {cnt}</span>'
                for nm, cnt in spk_counts.items()
            )
            st.markdown(badges, unsafe_allow_html=True)
    with col_b:
        notes = st.session_state.ss_slide_notes
        if notes:
            st.markdown(
                f'<div class="srow">'
                f'<span>📊 <strong>{len(notes)}</strong> slayt</span>'
                f'<span>🎭 <strong>{len(chars)}</strong> konuşmacı</span>'
                f'<span>📐 Slayt: <strong>{SLIDE_AREA_W}×{SLIDE_AREA_H}</strong></span>'
                f'</div>',
                unsafe_allow_html=True)
            for i, note in enumerate(notes[:6]):
                si  = min(st.session_state.ss_slide_speaker.get(i, 0), len(chars)-1)
                spk = chars[si]
                prev = (note[:88]+"…") if len(note)>88 else (note or "—")
                st.markdown(
                    f'<div class="sl-card" style="border-left-color:#{spk["hex"]};">'
                    f'<div class="sl-num" style="color:#{spk["hex"]};">'
                    f'{spk["emoji"]} {spk["name"]}  ·  Slayt {i+1}</div>'
                    f'<div class="sl-note">{prev}</div></div>',
                    unsafe_allow_html=True)
            if len(notes) > 6:
                st.caption(f"… ve {len(notes)-6} slayt daha")

    st.markdown("---")

    # ── ADIM 2: SES ───────────────────────────────────────────────────────────
    st.markdown('<p class="lbl">② Ses Dosyaları & Konuşmacı Atama</p>', unsafe_allow_html=True)
    if not st.session_state.ss_slide_notes:
        st.info("Önce bir PPTX dosyası yükleyin.")
    else:
        audio_assignment_ui(len(st.session_state.ss_slide_notes), chars)

    st.markdown("---")

    # ── ADIM 3: VİDEO ─────────────────────────────────────────────────────────
    st.markdown('<p class="lbl">③ Video Oluştur</p>', unsafe_allow_html=True)
    can_go = (
        st.session_state.ss_pptx_bytes is not None
        and PIL_OK and IMAGEIO_OK and PPTX_OK and LO_OK and PPM_OK and FFMPEG_OK
    )
    if not can_go:
        if not st.session_state.ss_pptx_bytes:
            st.info("Önce PPTX dosyası yükleyin.")
        else:
            missing = [n for n, ok in [
                ("Pillow",PIL_OK),("imageio",IMAGEIO_OK),("python-pptx",PPTX_OK),
                ("libreoffice",LO_OK),("poppler-utils",PPM_OK),("ffmpeg",FFMPEG_OK)
            ] if not ok]
            st.error(f"Eksik bağımlılık: {', '.join(missing)}")
            if not FFMPEG_OK:
                st.markdown(
                    '<div class="ffmpeg-warn">🔧 <b>ffmpeg:</b><br>'
                    'requirements.txt → <code>imageio[ffmpeg]</code><br>'
                    'packages.txt → <code>ffmpeg</code></div>',
                    unsafe_allow_html=True)
    else:
        n_slides   = len(st.session_state.ss_slide_notes)
        total_secs = sum(st.session_state.ss_durations.get(i, 3.0) for i in range(n_slides))
        total_with_jingles = total_secs + JINGLE_DUR * 2
        mins, secs = divmod(int(total_with_jingles), 60)
        st.markdown(
            f'<div class="srow">'
            f'<span>🎞️ <strong>{n_slides}</strong> slayt</span>'
            f'<span>⏱️ ~<strong>{mins}:{secs:02d}</strong> (jingle dahil)</span>'
            f'<span>📐 <strong>{VIDEO_W}×{VIDEO_H}</strong></span>'
            f'<span>🎬 <strong>{VIDEO_FPS} FPS</strong></span>'
            f'<span>🎭 <strong>{len(chars)}</strong> konuşmacı</span>'
            f'<span>🎵 5sn açılış + 5sn kapanış jingle</span>'
            f'</div>',
            unsafe_allow_html=True)
        c1, c2 = st.columns([3, 1])
        with c1:
            make_btn = st.button(
                "🎬 Video Oluştur", type="primary",
                use_container_width=True, key="wu_btn_make",
                disabled=(st.session_state.ss_video_bytes is not None))
        with c2:
            if st.button("🔄 Sıfırla", use_container_width=True, key="wu_btn_reset"):
                st.session_state.ss_video_bytes = None
                st.rerun()

        if make_btn:
            prog = st.progress(0)
            stat = st.empty()
            t0   = time.time()

            def cb(pct, msg):
                prog.progress(min(float(pct), 1.0))
                stat.markdown(
                    f"⚙️ **{msg}** &nbsp;"
                    f'<span style="color:#2a4038;font-size:.76rem;">— {time.time()-t0:.0f}s</span>',
                    unsafe_allow_html=True)

            work_dir = tempfile.mkdtemp(prefix="3soru_")
            try:
                cb(0.02, "Slaytlar görüntüye dönüştürülüyor…")
                slide_imgs = pptx_to_images(st.session_state.ss_pptx_bytes)
                n_actual   = len(slide_imgs)
                cb(0.10, f"Ses temizleniyor ve segmentler hazırlanıyor ({n_actual} slayt)…")
                audio_paths, seek_starts, dur_list = prepare_audio_segments(
                    slide_audio_map = st.session_state.ss_slide_audio,
                    durations       = st.session_state.ss_durations,
                    n_slides        = n_actual,
                    global_audio    = st.session_state.ss_global_audio,
                    use_global      = st.session_state.ss_use_global,
                    work_dir        = work_dir,
                )
                slide_speakers = [
                    chars[min(st.session_state.ss_slide_speaker.get(i, 0), len(chars)-1)]
                    for i in range(n_actual)
                ]
                video_bytes = build_video(
                    slide_images = slide_imgs,
                    audio_paths  = audio_paths,
                    seek_starts  = seek_starts,
                    durations    = dur_list,
                    speakers     = slide_speakers,
                    work_dir     = work_dir,
                    cb           = cb,
                )
                st.session_state.ss_video_bytes  = video_bytes
                st.session_state.ss_slide_images = slide_imgs
            except RuntimeError as e:
                st.error(f"❌ Hata:\n\n{e}")
            except Exception as e:
                import traceback
                st.error(f"❌ Beklenmedik hata: {e}")
                with st.expander("🔍 Traceback"):
                    st.code(traceback.format_exc())
            finally:
                shutil.rmtree(work_dir, ignore_errors=True)

        if st.session_state.ss_video_bytes:
            vb       = st.session_state.ss_video_bytes
            size     = len(vb)
            size_str = (f"{size//(1024*1024):.1f} MB"
                        if size > 1_048_576 else f"{size//1024:,} KB")
            st.success(f"✅ Video hazır — **{size_str}**")
            st.video(vb)
            col_d1, col_d2 = st.columns(2)
            with col_d1:
                st.download_button(
                    "⬇️ MP4 İndir", data=vb,
                    file_name="elif_aracıoglu_sunum.mp4", mime="video/mp4",
                    use_container_width=True, type="primary", key="wu_dl")
            with col_d2:
                b64v = base64.b64encode(vb).decode()
                st.markdown(
                    f'<a href="data:video/mp4;base64,{b64v}" '
                    'download="elif_aracıoglu_sunum.mp4" class="dl-alt">'
                    '📥 Alternatif İndirme</a>',
                    unsafe_allow_html=True)
            imgs = st.session_state.ss_slide_images
            if imgs:
                st.markdown("---")
                st.markdown('<p class="lbl">Slayt Önizlemeleri</p>', unsafe_allow_html=True)
                pc = st.columns(min(len(imgs), 4))
                for i, img in enumerate(imgs):
                    si  = min(st.session_state.ss_slide_speaker.get(i, 0), len(chars)-1)
                    spk = chars[si]
                    with pc[i % 4]:
                        st.image(img, caption=f"{spk['emoji']} {spk['name']} — S{i+1}",
                                 use_container_width=True)

if __name__ == "__main__":
    main()
